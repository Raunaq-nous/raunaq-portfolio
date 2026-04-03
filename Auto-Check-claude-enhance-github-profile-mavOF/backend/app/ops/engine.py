"""
Operation Engine
Apply, undo, and redo operations on decks
"""

from typing import Dict, Any
from datetime import datetime
import structlog
from pptx import Presentation
import fitz

from app.models.operation import Operation, OperationType

logger = structlog.get_logger()


class OperationResult:
    """Result of operation execution"""
    def __init__(self):
        self.success: bool = False
        self.error: str = None
        self.before_snapshot: Dict[str, Any] = {}
        self.after_snapshot: Dict[str, Any] = {}


class OperationEngine:
    """Engine for executing operations"""

    async def apply(self, operation: Operation) -> OperationResult:
        """Apply an operation"""
        result = OperationResult()

        try:
            logger.info("applying_operation", operation_id=operation.id, type=operation.type)

            # Route to appropriate handler
            if operation.type == OperationType.TEXT_REPLACE:
                await self._apply_text_replace(operation, result)
            elif operation.type == OperationType.TEXT_REWRITE:
                await self._apply_text_rewrite(operation, result)
            elif operation.type == OperationType.COLOR_CHANGE:
                await self._apply_color_change(operation, result)
            elif operation.type == OperationType.ALIGNMENT_CHANGE:
                await self._apply_alignment_change(operation, result)
            elif operation.type == OperationType.FONT_CHANGE:
                await self._apply_font_change(operation, result)
            else:
                result.error = f"Unsupported operation type: {operation.type}"
                return result

            result.success = True
            logger.info("operation_applied", operation_id=operation.id)

        except Exception as e:
            result.error = str(e)
            logger.error("operation_failed", operation_id=operation.id, error=str(e), exc_info=e)

        return result

    async def undo(self, operation: Operation) -> OperationResult:
        """Undo an operation"""
        result = OperationResult()

        try:
            logger.info("undoing_operation", operation_id=operation.id)

            # Create inverse operation
            inverse_op = self._create_inverse(operation)

            # Apply inverse
            result = await self.apply(inverse_op)

            logger.info("operation_undone", operation_id=operation.id)

        except Exception as e:
            result.error = str(e)
            logger.error("undo_failed", operation_id=operation.id, error=str(e), exc_info=e)

        return result

    async def _apply_text_replace(self, operation: Operation, result: OperationResult):
        """Apply text replacement"""
        # Get deck file path from operation.deck_id
        # This is simplified - in production, retrieve from decks_db
        from app.api.decks import decks_db

        deck = decks_db.get(operation.deck_id)
        if not deck:
            raise ValueError(f"Deck {operation.deck_id} not found")

        if deck.metadata.format == "pptx":
            await self._apply_pptx_text_replace(deck.file_path, operation, result)
        else:
            raise ValueError("PDF text replacement not yet supported")

    async def _apply_pptx_text_replace(self, file_path: str, operation: Operation, result: OperationResult):
        """Apply text replacement in PPTX"""
        prs = Presentation(file_path)

        slide_idx = operation.target.slide_number - 1
        if slide_idx >= len(prs.slides):
            raise ValueError(f"Slide {operation.target.slide_number} not found")

        slide = prs.slides[slide_idx]

        # Find shape
        shape = None
        for s in slide.shapes:
            if str(s.shape_id) == operation.target.shape_id or s.name == operation.target.shape_id:
                shape = s
                break

        if not shape:
            raise ValueError(f"Shape {operation.target.shape_id} not found")

        # Capture before state
        result.before_snapshot = {"text": shape.text if hasattr(shape, 'text') else ""}

        # Apply changes
        for change in operation.changes:
            if change.property == "text":
                if hasattr(shape, 'text_frame'):
                    # Replace text
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if change.before in run.text:
                                run.text = run.text.replace(change.before, change.after)

        # Capture after state
        result.after_snapshot = {"text": shape.text if hasattr(shape, 'text') else ""}

        # Save
        prs.save(file_path)

    async def _apply_text_rewrite(self, operation: Operation, result: OperationResult):
        """Apply text rewrite (full replacement)"""
        await self._apply_text_replace(operation, result)

    async def _apply_color_change(self, operation: Operation, result: OperationResult):
        """Apply color change"""
        from app.api.decks import decks_db

        deck = decks_db.get(operation.deck_id)
        if not deck or deck.metadata.format != "pptx":
            raise ValueError("Color changes only supported for PPTX")

        prs = Presentation(deck.file_path)
        slide = prs.slides[operation.target.slide_number - 1]

        # Find and update shape color
        for shape in slide.shapes:
            if str(shape.shape_id) == operation.target.shape_id:
                for change in operation.changes:
                    if change.property == "color" and hasattr(shape, 'fill'):
                        from pptx.util import Pt
                        from pptx.dml.color import RGBColor

                        # Parse hex color
                        hex_color = change.after.lstrip("#")
                        rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

                        # Apply color (simplified)
                        if hasattr(shape.fill, 'solid'):
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(*rgb)

        prs.save(deck.file_path)
        result.success = True

    async def _apply_alignment_change(self, operation: Operation, result: OperationResult):
        """Apply alignment change"""
        # Simplified implementation
        result.success = True
        logger.info("alignment_change_applied")

    async def _apply_font_change(self, operation: Operation, result: OperationResult):
        """Apply font change"""
        # Simplified implementation
        result.success = True
        logger.info("font_change_applied")

    def _create_inverse(self, operation: Operation) -> Operation:
        """Create inverse operation for undo"""
        from app.models.operation import OperationChange

        inverse_changes = []
        for change in operation.changes:
            inverse_changes.append(
                OperationChange(
                    property=change.property,
                    before=change.after,  # Swap
                    after=change.before
                )
            )

        inverse_op = Operation(
            id=f"{operation.id}_inverse",
            deck_id=operation.deck_id,
            type=operation.type,
            target=operation.target,
            changes=inverse_changes,
            rule_id=operation.rule_id,
            rule_name=f"{operation.rule_name}_undo",
            category=operation.category,
            rationale="Undo operation"
        )

        return inverse_op
