"""
PPTX Parser
Parse PowerPoint files and extract asset graph
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from typing import List, Dict, Optional, Any
import structlog
from datetime import datetime

from app.models.deck import DeckMetadata
from app.models.asset import (
    Asset, AssetType, AssetGraph, BoundingBox,
    ColorInfo, TextRun, ChartData, TableData
)

logger = structlog.get_logger()


class PPTXParser:
    """Parser for PPTX files"""

    async def extract_metadata(self, file_path: str) -> DeckMetadata:
        """Extract metadata from PPTX file"""
        try:
            prs = Presentation(file_path)

            # Get core properties
            core_props = prs.core_properties

            metadata = DeckMetadata(
                title=core_props.title or None,
                author=core_props.author or None,
                created_at=core_props.created,
                modified_at=core_props.modified,
                slide_count=len(prs.slides),
                format="pptx",
                file_size=0,  # Will be set by caller
                page_dimensions={
                    "width": prs.slide_width.inches if prs.slide_width else 10,
                    "height": prs.slide_height.inches if prs.slide_height else 7.5
                }
            )

            logger.info(
                "pptx_metadata_extracted",
                slides=metadata.slide_count,
                title=metadata.title
            )

            return metadata

        except Exception as e:
            logger.error("pptx_metadata_failed", file_path=file_path, error=str(e))
            raise

    async def parse_to_asset_graph(self, file_path: str, deck_id: str) -> AssetGraph:
        """Parse PPTX and create asset graph"""
        try:
            prs = Presentation(file_path)
            asset_graph = AssetGraph(deck_id=deck_id)

            asset_counter = 0

            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_assets = []

                for shape in slide.shapes:
                    asset = await self._parse_shape(shape, slide_idx, asset_counter)
                    if asset:
                        slide_assets.append(asset)
                        asset_graph.asset_index[asset.id] = asset
                        asset_counter += 1

                asset_graph.slides[slide_idx] = slide_assets

            # Calculate metrics
            asset_graph.total_assets = len(asset_graph.asset_index)
            asset_graph.assets_by_type = {}
            for asset in asset_graph.asset_index.values():
                type_name = asset.type.value
                asset_graph.assets_by_type[type_name] = \
                    asset_graph.assets_by_type.get(type_name, 0) + 1

            logger.info(
                "pptx_parsed",
                deck_id=deck_id,
                slides=len(prs.slides),
                assets=asset_graph.total_assets
            )

            return asset_graph

        except Exception as e:
            logger.error("pptx_parse_failed", file_path=file_path, error=str(e))
            raise

    async def _parse_shape(
        self, shape, slide_number: int, asset_id: int
    ) -> Optional[Asset]:
        """Parse individual shape"""
        try:
            # Determine asset type
            asset_type = self._get_asset_type(shape)
            if not asset_type:
                return None

            # Extract bounding box
            bbox = BoundingBox(
                x=shape.left.inches if hasattr(shape, 'left') else 0,
                y=shape.top.inches if hasattr(shape, 'top') else 0,
                width=shape.width.inches if hasattr(shape, 'width') else 0,
                height=shape.height.inches if hasattr(shape, 'height') else 0
            )

            # Extract text content
            text_content = None
            text_runs = []
            if hasattr(shape, 'text'):
                text_content = shape.text
                if hasattr(shape, 'text_frame'):
                    text_runs = self._extract_text_runs(shape.text_frame)

            # Extract colors
            fill_color = self._extract_fill_color(shape)
            line_color = self._extract_line_color(shape)

            # Extract chart data if chart
            chart_data = None
            if asset_type == AssetType.CHART and hasattr(shape, 'chart'):
                chart_data = self._extract_chart_data(shape.chart)

            # Extract table data if table
            table_data = None
            if asset_type == AssetType.TABLE and hasattr(shape, 'table'):
                table_data = self._extract_table_data(shape.table)

            asset = Asset(
                id=f"asset_{asset_id}",
                slide_number=slide_number,
                shape_id=str(shape.shape_id) if hasattr(shape, 'shape_id') else f"shape_{asset_id}",
                shape_name=shape.name if hasattr(shape, 'name') else None,
                type=asset_type,
                bbox=bbox,
                z_order=0,  # TODO: Extract z-order
                rotation=shape.rotation if hasattr(shape, 'rotation') else 0,
                fill_color=fill_color,
                line_color=line_color,
                text_content=text_content,
                text_runs=text_runs,
                chart_data=chart_data,
                table_data=table_data
            )

            return asset

        except Exception as e:
            logger.warning("shape_parse_failed", error=str(e))
            return None

    def _get_asset_type(self, shape) -> Optional[AssetType]:
        """Determine asset type from shape"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if not hasattr(shape, 'shape_type'):
            return None

        shape_type = shape.shape_type

        if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return AssetType.TEXT_BOX
        elif shape_type == MSO_SHAPE_TYPE.CHART:
            return AssetType.CHART
        elif shape_type == MSO_SHAPE_TYPE.TABLE:
            return AssetType.TABLE
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            return AssetType.IMAGE
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            return AssetType.GROUP
        elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return AssetType.SHAPE
        elif shape_type == MSO_SHAPE_TYPE.LINE or shape_type == MSO_SHAPE_TYPE.CONNECTOR:
            return AssetType.CONNECTOR
        else:
            return AssetType.SHAPE

    def _extract_text_runs(self, text_frame) -> List[TextRun]:
        """Extract text runs with formatting"""
        runs = []
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    text_run = TextRun(
                        text=run.text,
                        font_name=run.font.name if run.font.name else None,
                        font_size=run.font.size.pt if run.font.size else None,
                        bold=run.font.bold or False,
                        italic=run.font.italic or False,
                        underline=run.font.underline or False
                    )
                    runs.append(text_run)
        except Exception as e:
            logger.warning("text_runs_extract_failed", error=str(e))

        return runs

    def _extract_fill_color(self, shape) -> Optional[ColorInfo]:
        """Extract fill color"""
        try:
            if hasattr(shape, 'fill') and shape.fill.type:
                from pptx.enum.dml import MSO_FILL_TYPE
                if shape.fill.type == MSO_FILL_TYPE.SOLID:
                    color = shape.fill.fore_color
                    if hasattr(color, 'rgb'):
                        rgb = color.rgb
                        return ColorInfo(
                            rgb=(rgb[0], rgb[1], rgb[2]),
                            hex=f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                        )
        except Exception as e:
            logger.debug("fill_color_extract_failed", error=str(e))

        return None

    def _extract_line_color(self, shape) -> Optional[ColorInfo]:
        """Extract line color"""
        try:
            if hasattr(shape, 'line') and shape.line.fill.type:
                from pptx.enum.dml import MSO_FILL_TYPE
                if shape.line.fill.type == MSO_FILL_TYPE.SOLID:
                    color = shape.line.fill.fore_color
                    if hasattr(color, 'rgb'):
                        rgb = color.rgb
                        return ColorInfo(
                            rgb=(rgb[0], rgb[1], rgb[2]),
                            hex=f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                        )
        except Exception as e:
            logger.debug("line_color_extract_failed", error=str(e))

        return None

    def _extract_chart_data(self, chart) -> Optional[ChartData]:
        """Extract chart data"""
        try:
            series_data = []
            for series in chart.series:
                series_data.append({
                    "name": series.name,
                    "values": list(series.values)
                })

            categories = []
            if hasattr(chart, 'plots') and len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, 'categories'):
                    categories = [str(cat) for cat in plot.categories]

            return ChartData(
                chart_type=str(chart.chart_type) if hasattr(chart, 'chart_type') else "unknown",
                series=series_data,
                categories=categories,
                axes={},
                legend={}
            )

        except Exception as e:
            logger.warning("chart_data_extract_failed", error=str(e))
            return None

    def _extract_table_data(self, table) -> Optional[TableData]:
        """Extract table data"""
        try:
            cells = []
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    row_cells.append(cell.text if hasattr(cell, 'text') else "")
                cells.append(row_cells)

            return TableData(
                rows=len(table.rows),
                cols=len(table.columns) if hasattr(table, 'columns') else 0,
                cells=cells,
                header_row=False,  # TODO: Detect header
                header_col=False
            )

        except Exception as e:
            logger.warning("table_data_extract_failed", error=str(e))
            return None
